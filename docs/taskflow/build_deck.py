"""Rework the Round 1 deck into the Round 2 business proposal.

Keeps the Accenture master, theme and 16:9 geometry by opening the Round 1 file and
adding slides to it. Slides 1-4 (cover, team, problem, solution flow) are kept: the
Round 1 solution slide already promised "AI proposes the checks that would confirm or
kill each candidate", which Round 2 delivers -- so it reads as continuity, not filler.
"""
import copy, pathlib, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

SRC = "details/AIC_Round1_LedgerLens.pptx"
OUT = "submission/LedgerLens_Business_Proposal.pptx"
SHOTS = pathlib.Path("docs/screenshots")
PURPLE = RGBColor(0x5B, 0x21, 0xB6)
INK = RGBColor(0x1A, 0x1D, 0x21)
MUTE = RGBColor(0x55, 0x5A, 0x63)

prs = Presentation(SRC)
W, H = prs.slide_width, prs.slide_height
L_TEXT = prs.slide_masters[4].slide_layouts[16]   # Content 1: light mode  (title + body)
L_TITLE = prs.slide_masters[4].slide_layouts[12]  # Content: title only
L_KEY = prs.slide_masters[4].slide_layouts[37]    # Key message: light mode

def _title(slide, text, size=27):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = PURPLE
            return ph
    return None

def _clear_unused(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 0 and not ph.text_frame.text.strip():
            ph._element.getparent().remove(ph._element)

def text_slide(title, blocks, sub=None):
    """blocks: list of (bullet_text, indent_level, bold)"""
    s = prs.slides.add_slide(L_TEXT)
    _title(s, title)
    body = None
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            body = ph; break
    if body is None:
        body = s.shapes.add_textbox(Inches(0.9), Inches(1.7), W - Inches(1.8), H - Inches(2.4))
    tf = body.text_frame; tf.clear(); tf.word_wrap = True
    first = True
    if sub:
        p = tf.paragraphs[0]; p.text = sub; first = False
        for r in p.runs:
            r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = MUTE
    for text, lvl, bold in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text; p.level = lvl
        p.space_after = Pt(7)
        for r in p.runs:
            r.font.size = Pt(15 if lvl == 0 else 12.5)
            r.font.bold = bold
            r.font.color.rgb = INK if lvl == 0 else MUTE
    _clear_unused(s)
    return s

def shot_slide(title, image, caption, crop_left=0.0, crop_bottom=1.0):
    s = prs.slides.add_slide(L_TITLE)
    _title(s, title, size=24)
    img = SHOTS / image
    im = Image.open(img).convert("RGB")
    # Crop the sidebar away. On a 13.3in slide the app's full 1500px width renders the
    # card text below ~9pt, which is unreadable from the back of a room; dropping the
    # chrome roughly doubles the effective size of the thing being pointed at.
    if crop_left or crop_bottom < 1.0:
        im = im.crop((int(im.width * crop_left), 0, im.width, int(im.height * crop_bottom)))
    # Downscale: a 3000px screenshot is ~400 KB and the base deck is already 19 MB.
    scale = 1700 / im.width
    im = im.resize((1700, int(im.height * scale)), Image.LANCZOS)
    buf = io.BytesIO(); im.save(buf, "JPEG", quality=82, optimize=True); buf.seek(0)
    top = Inches(1.55)
    avail_w, avail_h = W - Inches(1.4), H - top - Inches(1.05)
    ar = im.width / im.height
    w = avail_w; h = int(w / ar)
    if h > avail_h:
        h = avail_h; w = int(h * ar)
    s.shapes.add_picture(buf, int((W - w) / 2), top, w, h)
    tb = s.shapes.add_textbox(Inches(0.7), H - Inches(0.92), W - Inches(1.4), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; tf.text = caption
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(12); r.font.color.rgb = MUTE
    _clear_unused(s)
    return s

def key_slide(text):
    s = prs.slides.add_slide(L_KEY)
    _title(s, text, size=30)
    _clear_unused(s)
    return s

# ----------------------------------------------------------------- the slides

key_slide("Round 1 promised an AI investigator that proposes the checks\n"
          "which confirm or kill each candidate.\n\nRound 2 ships exactly that.")

text_slide("What we built", [
    ("A KPI intelligence-to-action engine: root-cause analysis as a set intersection over a ledger of business changes, verified by negative controls.", 0, False),
    ("Diagnoses a −$416,144 renewals shortfall in 1.3 seconds — and rejects the plausible-but-wrong cause a human under time pressure would have acted on.", 0, True),
    ("307 passing tests. Runs end to end with no API key: the ranking path never calls a model.", 1, False),
    ("22 distinct query_ids replayable on a single diagnosis card.", 1, False),
    ("Ten of ten Minimum Prototype Expectations close. Eight of eight Round 2 objectives.", 1, False),
], sub="github.com/omega-sus67/LedgerLens")

text_slide("How it decides — five components, every one a query", [
    ("score = 0.25·T + 0.30·C + 0.15·D + 0.25·N + 0.05·P", 0, True),
    ("T  temporal — the change began shortly before the metric broke. Precedence is necessary, never sufficient.", 1, False),
    ("C  cohort match — row-level Jaccard of blast radius against the affected cohort. A wide radius scores badly even if it was the cause.", 1, False),
    ("D  dose–response — rank correlation of exposure against impact. Uninformative here, and the card says so rather than manufacturing a number.", 1, False),
    ("N  negative controls — the fraction of falsifiable predictions that survived. A passing control is a failure to falsify, not a confirmation.", 1, False),
    ("P  learned prior — Beta–Bernoulli over analyst verdicts. Weighted 0.05: it sharpens a ranking, it never decides one.", 1, False),
    ("No component is an LLM opinion. All five are reproducible queries — that is the answer to \"how is this not hallucination?\"", 0, True),
])

shot_slide("The mechanism that matters: the decoy dies",
           "04_decoy_rejected.png",
           "Marketing cut spend in DACH one day before renewals broke — temporally more plausible than the truth. "
           "The engine predicts DACH Mid/SMB should also have dropped. They came in flat at −1.31%. "
           "That control failure is decisive: rejected at 0.322, not merely outranked.",
           crop_left=0.245, crop_bottom=0.80)

text_slide("LLM versus non-LLM — the split, stated exactly", [
    ("\"The LLM should not be treated as the source of quantitative truth.\"  — Round 2 brief", 0, True),
    ("Non-LLM — detection, drill-down, cohort intersection, all five score components, every negative control, the rejection of the decoy. 89 SQL queries, no model involved.", 0, False),
    ("LLM — proposing additional checks (which this engine then executes in SQL), naming causes the connected data cannot test, and writing the prose.", 0, False),
    ("The boundary is enforced in code, not by convention: proposed checks are built with decisive=False — the field the N-scorer reads — and are never passed to it.", 0, True),
    ("A test asserts every score is byte-identical with the AI lane on and off.", 1, False),
])

text_slide("The AI investigator lane — three call sites, all additive", [
    ("1 · Proposed checks — the model fills a fixed template vocabulary; it never writes SQL. We execute. Hallucinated dimensions are rejected before becoming a query, and the count is shown: \"4 accepted, 2 rejected by validation.\"", 0, False),
    ("2 · Unverifiable causes — a separately labelled panel of explanations the connected data cannot test, each naming the feed that would settle it.", 0, False),
    ("3 · Guarded narration — every numeric token in the model's prose must appear in the verified payload. One invented digit discards the narration for the deterministic template, and the page says so.", 0, False),
    ("Provider-agnostic by construction: Gemini 2.5 Flash by default, Claude Sonnet via one environment variable. Two adapters sharing no code path, behind a one-method interface.", 0, True),
    ("~$0.0048 per diagnosis with the lane on. $0.0000 with it off.", 1, False),
])

text_slide("Target users — four audiences, one computation", [
    ("Persona is accepted only by the narrator, downstream of every query — so it cannot reach a query. The evidence behind all four cards is identical by construction.", 0, False),
    ("Revenue Analyst — every control, every query id, the full drill-down lattice.", 1, False),
    ("CFO — dollars and forecast risk, and an escalation. Never an instruction to roll back a release: she does not hold that lever.", 1, False),
    ("Payments On-Call — the event id, the blast radius, the rollback.", 1, False),
    ("Growth Marketing — their own KPI's real story, and a named policy where a cut is withheld.", 1, False),
    ("Decision rights are mechanical: a persona that does not hold a lever is shown an escalation, never an instruction.", 0, True),
])

shot_slide("Evidence, telemetry and cost — shown, not asserted",
           "09_telemetry.png",
           "Every number on the page came from a logged SQL query with a replayable query_id. "
           "\"Queries\" is three different numbers and the panel refuses to merge them: what the diagnosis cost, "
           "what a reader can audit, and what the cache avoided.",
           crop_left=0.245)

text_slide("Business case and impact", [
    ("Stated assumptions — mid-market B2B SaaS at ~€50M ARR: ~40 material KPI movements a year, ~3 analyst-days each, ~$460/day fully loaded. Directional sizing, not measured results.", 0, False),
    ("Analyst time recovered — ~$55k/year of capacity redirected from evidence assembly to judgement. Real, but the smaller half.", 0, False),
    ("Wrong-action avoidance — the larger half. Acting on the decoy leaves the −$410k leak running AND reverses a campaign on a false premise. One avoided misattribution dominates the entire labour saving.", 0, True),
    ("Decision-window preservation — 1.3 seconds instead of three days means the finding lands while the fix is still cheap.", 0, False),
    ("Cost to run — one DuckDB file. No graph database, no vector store, no agent framework. Under $1/year of model spend at 40 diagnoses.", 0, False),
])

text_slide("Phased roadmap", [
    ("Shipped — Round 2 prototype", 0, True),
    ("3 KPIs across 3 sources · semantic contract · drill-down · 5-component scoring · 5 control rules · 4 personas · entitlement · abstention · AI investigator lane · feedback loop · telemetry.", 1, False),
    ("Phase 1 — Rigour (1–2 months)", 0, True),
    ("Difference-in-differences with bootstrap CI · calendar-regressor baseline enabling bidirectional detection · the discriminating test between near-tied candidates · a bounded exploration pass before abstaining.", 1, False),
    ("Phase 2 — Enterprise integration (3–6 months)", 0, True),
    ("Real connectors replacing synthetic fixtures — the mapping is the only thing that changes. Warehouse-native execution on Snowflake/Databricks/BigQuery, since the engine is already just SQL. SSO and per-user verdict attribution. Push delivery to Slack, email, PagerDuty.", 1, False),
    ("Phase 3 — Proactive (6–12 months)", 0, True),
    ("A watchtower that surfaces anomalies rather than waiting to be pointed · the LLM event normalizer, gated behind confidence calibration · cross-KPI interaction detection.", 1, False),
])

text_slide("Key risks and mitigations", [
    ("The system names a confident wrong cause → negative controls try to kill each candidate; a decisive failure rejects outright and cannot be outvoted.  SHIPPED", 0, False),
    ("The LLM hallucinates a number → the numbers guard discards the whole narration for the deterministic template.  SHIPPED", 0, False),
    ("The LLM invents a dimension or metric → validation against the registry before execution; rejections counted and displayed.  SHIPPED", 0, False),
    ("The LLM quietly influences the verdict → proposed checks built decisive=False, never passed to the scorer; asserted by test.  SHIPPED", 0, False),
    ("Blast-radius metadata is wrong → the failure direction is the mitigation: too wide fails its controls, too narrow leaves nothing above the floor. It degrades toward \"I don't know\".  SHIPPED", 0, False),
    ("Feedback poisons the ranking → P weighted 0.05; a prior saturated with 200 confirmations still cannot clear the floor or rescue what a control killed.  SHIPPED", 0, False),
    ("Vendor lock-in or outage → one provider table, two live adapters; every failure returns a recorded reason, so \"found nothing\" never looks like \"vendor was down\".  SHIPPED", 0, False),
    ("Verdicts unattributable to a person → no auth in this build. Named, not fixed: SSO in Phase 2.", 0, False),
])

text_slide("How this maps to the rubric", [
    ("Ten of ten Minimum Prototype Expectations close.", 0, True),
    ("3 KPIs / 3 sources, different grains · semantic contract · 4 personas · multi-factor movement decomposed · two low-confidence scenarios · sparse-history KPI · role entitlement · evidence with freshness and lineage · LLM vs non-LLM breakdown · runtime telemetry.", 1, False),
    ("Eight of eight Round 2 objectives.", 0, True),
    ("Detection and prioritisation · heterogeneous reconciliation · ranked drivers · persona narratives with traceable evidence · uncertainty and abstention · actions grounded in levers and decision rights · a learning mechanism from analyst feedback · realistic security, cost and latency.", 1, False),
    ("What is deliberately not built is named too — the LLM event normalizer is cut because it is the one design that puts a model on the ranking path.", 0, True),
])

# ------------------------------------------------------- reorder + finalise
#
# python-pptx appends; the Video and Thank-you slides must stay last. Move them to the
# end of the slide-id list rather than rebuilding the deck.
ids = prs.slides._sldIdLst
kids = list(ids)
video_i, thanks_i = 4, 5          # 0-based positions in the Round 1 deck
tail = [kids[video_i], kids[thanks_i]]
for el in tail:
    ids.remove(el)
for el in tail:
    ids.append(el)

# Fill the two Round 1 placeholders.
TEAM_NAME = pathlib.Path("submission/.team_name").read_text().strip() if pathlib.Path("submission/.team_name").exists() else None
VIDEO_URL = pathlib.Path("submission/.video_url").read_text().strip() if pathlib.Path("submission/.video_url").exists() else None
filled = []

def _frames(slide):
    """Every text frame on a slide, INCLUDING table cells -- the Round 1 team-name
    placeholder lives in a table, which has no text_frame of its own."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            yield sh.text_frame
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    yield cell.text_frame

for s in prs.slides:
    for tf in _frames(s):
        for para in tf.paragraphs:
            for r in para.runs:
                if TEAM_NAME and "[ADD TEAM NAME]" in r.text:
                    r.text = r.text.replace("[ADD TEAM NAME]", TEAM_NAME); filled.append("team name")
                if VIDEO_URL and "[ADD VIDEO LINK HERE" in r.text:
                    r.text = VIDEO_URL; filled.append("video url")

pathlib.Path("submission").mkdir(exist_ok=True)
TMP_OUT = OUT + ".tmp"
prs.save(TMP_OUT)

# ---- shrink the cover art.
# The Round 1 template ships its cover background at 9000x5209 -- 12 MB, two thirds of
# the whole deck. A 13.3in slide at 200 dpi needs ~2670px; the rest is invisible weight
# on an upload form with a size cap. Rewritten in place so every rel stays valid.
import zipfile, io
zin = zipfile.ZipFile(TMP_OUT)
BIG = "ppt/media/image30.jpg"
rewritten = {}
if BIG in zin.namelist():
    im = Image.open(io.BytesIO(zin.read(BIG))).convert("RGB")
    if im.width > 2600:
        im = im.resize((2600, round(im.height * 2600 / im.width)), Image.LANCZOS)
    b = io.BytesIO(); im.save(b, "JPEG", quality=84, optimize=True, progressive=True)
    rewritten[BIG] = b.getvalue()
with zipfile.ZipFile(OUT, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        zout.writestr(item, rewritten.get(item.filename) or zin.read(item.filename))
zin.close()
pathlib.Path(TMP_OUT).unlink()
size = pathlib.Path(OUT).stat().st_size
print(f"wrote {OUT}: {len(prs.slides.__iter__.__self__._sldIdLst)} slides, {size/1048576:.1f} MB")
print("filled:", filled or "(nothing -- placeholders remain)")
remaining = []
for i, s in enumerate(prs.slides, 1):
    for tf in _frames(s):
        if "[ADD" in tf.text:
            remaining.append((i, tf.text[:60].replace("\n", " / ")))
print("REMAINING PLACEHOLDERS:", remaining or "none")
